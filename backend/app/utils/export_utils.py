# app/utils/export_utils.py
from docx import Document
from fastapi.responses import StreamingResponse
import io
from pptx import Presentation
from pptx.util import Pt

def assemble_docx(project, sections):
    doc = Document()
    doc.add_heading(project.name or "Document", level=1)
    if project.main_topic:
        doc.add_paragraph(f"Topic: {project.main_topic}")
    for sec in sorted(sections, key=lambda s: s.order):
        doc.add_heading(sec.title, level=2)
        p = doc.add_paragraph(sec.content or "")
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    filename = f"{project.name or 'project'}".replace(" ", "_") + ".docx"
    return StreamingResponse(buf,
                             media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                             headers={"Content-Disposition": f"attachment; filename={filename}"})

def assemble_pptx(project, sections):
    prs = Presentation()
    title_slide = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_slide)
    s.shapes.title.text = project.name or "Presentation"
    subtitle = s.placeholders[1]
    subtitle.text = project.main_topic or ""
    for sec in sorted(sections, key=lambda s: s.order):
        layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sec.title
        tf = slide.shapes.placeholders[1].text_frame
        # break content into paragraphs
        text = sec.content or ""
        for para in text.split("\n"):
            p = tf.add_paragraph()
            p.text = para
            p.level = 0
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = f"{project.name or 'presentation'}".replace(" ", "_") + ".pptx"
    return StreamingResponse(buf,
                             media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                             headers={"Content-Disposition": f"attachment; filename={filename}"})
